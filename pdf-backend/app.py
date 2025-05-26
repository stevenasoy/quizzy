from flask import Flask, request, jsonify
from flask_cors import CORS
import io
from pptx import Presentation
from docx import Document
import pdfplumber
import os
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

app = Flask(__name__)
CORS(app)

def get_file_type(filename, content_type):
    """Determine file type from both filename and content type"""
    ext = filename.lower().split('.')[-1] if '.' in filename else ''
    if ext in ['doc', 'docx'] or 'word' in content_type:
        return 'docx'
    elif ext in ['ppt', 'pptx'] or 'powerpoint' in content_type:
        return 'pptx'
    elif ext == 'pdf' or 'pdf' in content_type:
        return 'pdf'
    elif ext == 'txt' or 'text/plain' in content_type:
        return 'txt'
    return None

def extract_text_from_pdf(file_stream):
    """Extract text from PDF file"""
    try:
        with pdfplumber.open(file_stream) as pdf:
            text = ''
            for page in pdf.pages:
                text += page.extract_text() or ''
        return text.strip()
    except Exception as e:
        return f"Error extracting PDF text: {str(e)}"

def extract_text_from_pptx(file_stream):
    """Extract text from PPTX file"""
    try:
        prs = Presentation(file_stream)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text).strip()
    except Exception as e:
        return f"Error extracting PPTX text: {str(e)}"

def extract_text_from_docx(file_stream):
    """Extract text from DOCX file"""
    try:
        doc = Document(file_stream)
        text = []
        for paragraph in doc.paragraphs:
            text.append(paragraph.text)
        return '\n'.join(text).strip()
    except Exception as e:
        return f"Error extracting DOCX text: {str(e)}"

@app.route('/extract-file', methods=['POST'])
def extract_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400

    try:
        file_type = get_file_type(file.filename, file.content_type)
        if not file_type:
            return jsonify({'error': 'Unsupported file type'}), 400

        file_stream = io.BytesIO(file.read())
        
        if file_type == 'pdf':
            text = extract_text_from_pdf(file_stream)
        elif file_type == 'pptx':
            text = extract_text_from_pptx(file_stream)
        elif file_type == 'docx':
            text = extract_text_from_docx(file_stream)
        elif file_type == 'txt':
            text = file_stream.read().decode('utf-8')
        else:
            return jsonify({'error': 'Unsupported file type'}), 400

        if not text:
            return jsonify({'error': 'No text could be extracted'}), 400

        return jsonify({'text': text})

    except Exception as e:
        return jsonify({'error': f'Error processing file: {str(e)}'}), 500

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5001, debug=True) 